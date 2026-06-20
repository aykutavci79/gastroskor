"""GastroSkor üye restoran panel onboarding sunumu (.pptx) — e-posta eki."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("GastroSkor-Panel-Uye-Onboarding.pptx")

# GastroSkor koyu tema
BG = RGBColor(0x14, 0x14, 0x14)
PANEL = RGBColor(0x1E, 0x1E, 0x1E)
TEXT = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xA0, 0xA0, 0xA0)
ACCENT = RGBColor(0xFF, 0xB7, 0x03)
SUCCESS = RGBColor(0x52, 0xB7, 0x88)
BORDER = RGBColor(0x2E, 0x2E, 0x2E)


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def _add_bullets(slide, left, top, width, height, lines: list[str], *, size: int = 16) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(8)
        if p.runs:
            run = p.runs[0]
        else:
            run = p.add_run()
            run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = TEXT
        run.font.name = "Calibri"


def _slide_title(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_textbox(slide, Inches(0.7), Inches(0.55), Inches(8.6), Inches(0.9), title, size=30, bold=True, color=ACCENT)
    if subtitle:
        _add_textbox(slide, Inches(0.7), Inches(1.35), Inches(8.6), Inches(0.6), subtitle, size=14, color=MUTED)


def _slide_content(prs: Presentation, title: str, bullets: list[str], footer: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_textbox(slide, Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.7), title, size=26, bold=True, color=TEXT)
    _add_bullets(slide, Inches(0.85), Inches(1.35), Inches(8.4), Inches(5.2), bullets, size=17)
    if footer:
        _add_textbox(slide, Inches(0.7), Inches(6.55), Inches(8.6), Inches(0.5), footer, size=12, color=MUTED)


def _slide_pricing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_textbox(slide, Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.7), "Üyelik ve ücretlendirme", size=26, bold=True)

    rows = [
        ("1. ay", "Deneme", "0 TL", "Tüm panel özellikleri açık"),
        ("2. ay", "Giriş", "399 TL / ay", "Deneme sonrası ilk ücretli dönem"),
        ("3. ay ve sonrası", "Standart", "599 TL / ay", "Tam panel + görünürlük + kampanyalar"),
    ]
    y = Inches(1.5)
    for period, label, price, note in rows:
        shape = slide.shapes.add_shape(1, Inches(0.8), y, Inches(8.4), Inches(1.15))  # MSO_SHAPE.RECTANGLE
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = BORDER
        _add_textbox(slide, Inches(1.0), y + Inches(0.12), Inches(2.2), Inches(0.35), period, size=14, bold=True, color=ACCENT)
        _add_textbox(slide, Inches(3.2), y + Inches(0.12), Inches(1.5), Inches(0.35), label, size=14, color=MUTED)
        _add_textbox(slide, Inches(4.8), y + Inches(0.12), Inches(1.8), Inches(0.35), price, size=20, bold=True, color=SUCCESS)
        _add_textbox(slide, Inches(1.0), y + Inches(0.55), Inches(7.8), Inches(0.45), note, size=13, color=TEXT)
        y += Inches(1.35)

    _add_bullets(
        slide,
        Inches(0.85),
        Inches(5.35),
        Inches(8.2),
        Inches(1.5),
        [
            "Ödeme panel üzerinden (kart) — fatura bilgilerinizi güncel tutun.",
            "Deneme bitiminden önce imzalı sözleşmeyi posta ile göndermeniz gerekir.",
            "Ek AI analiz paketleri panelde ayrıca listelenir (isteğe bağlı).",
        ],
        size=14,
    )


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Kapak
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BG)
    _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.0), "GastroSkor", size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_textbox(
        slide,
        Inches(0.8),
        Inches(3.0),
        Inches(8.4),
        Inches(0.8),
        "Restoran Paneli — Üyelik Aktivasyon Kılavuzu",
        size=24,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(1.2),
        Inches(4.0),
        Inches(7.6),
        Inches(1.2),
        "Bu dosya, üyeliğiniz aktive edildiğinde e-posta ile gönderilir.\n"
        "Paneli kendi hızınızda inceleyin; her bölümde ne işe yaradığını ve nasıl kullanılacağını bulacaksınız.",
        size=15,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(slide, Inches(0.8), Inches(6.2), Inches(8.4), Inches(0.4), "www.gastroskor.com.tr/panel", size=14, color=SUCCESS, align=PP_ALIGN.CENTER)

    _slide_content(
        prs,
        "GastroSkor paneli size ne kazandırır?",
        [
            "Mekanınızı mobil uygulama ve web’de öne çıkarır — menü, fotoğraf, online sipariş.",
            "Müşteri geri bildirimini yönetirsiniz: kötü yorum yayına çıkmadan önce çözüm şansı.",
            "Takipçilerinize toplu bildirim ve kupon kampanyası gönderebilirsiniz.",
            "Rakip restoranları AI ile kıyaslayıp güçlü / zayıf yönlerinizi görürsünüz.",
            "Online siparişleri tek ekrandan onaylar, reddeder veya yazdırırsınız.",
            "Google İşletme Profilinizi bağlayarak tüm Google yorum havuzunu analiz edebilirsiniz.",
        ],
        footer="Panel web üzerinden çalışır — bilgisayar veya tablet tarayıcısı yeterli.",
    )

    _slide_content(
        prs,
        "Panele nasıl girilir?",
        [
            "Adres: www.gastroskor.com.tr/panel",
            "Google işletme e-postanız ile giriş yapın (başvuruda kullandığınız hesap).",
            "İlk girişte mekanınızı doğrulayın: SMS kodu veya ziyaret doğrulaması.",
            "Tam yetki (mesaj, kupon, remedy teklifi) için doğrulama tamamlanmalıdır.",
            "Deneme süresince üst bantta kalan gün ve sonraki ay ücreti görünür.",
        ],
        footer="Mobil uygulamada panel sekmesi yok — tarayıcıdan kullanın.",
    )

    _slide_content(
        prs,
        "Ana sayfa (Dashboard)",
        [
            "Haftalık harita tıklama ve profil görüntülenme sayıları.",
            "Google ve GastroSkor puan özeti.",
            "Açık şikayet ve bekleyen online sipariş sayısı.",
            "AI özet metni: performans ve rakip kısa değerlendirme.",
            "Rakip ekleme ve analiz başlatma kısayolları.",
        ],
    )

    _slide_content(
        prs,
        "Vitrin, menü ve görünürlük",
        [
            "Menü fotoğrafı ve kart kapak görseli yükleyin.",
            "Menü kalemlerini ekleyin: ad, fiyat, kategori, açıklama.",
            "Online sipariş aç/kapa, mutfak etiketleri, kurye rozeti.",
            "Sipariş telefonu, WhatsApp ve Instagram bağlantıları.",
            "Sesli sipariş için ürün kataloğunu işaretleyin (online sipariş açıkken).",
            "Ayarlar kaydedildiğinde uygulamada ve web’de vitrin güncellenir.",
        ],
        footer="Menü ve vitrin, aktif üyelik / deneme süresince yayında kalır.",
    )

    _slide_content(
        prs,
        "Online siparişler",
        [
            "Müşteri uygulamadan sipariş verir → panelde anlık bildirim (ses + tarayıcı).",
            "Sipariş kartında: ürünler, adres, telefon, not.",
            "Onayla veya reddet (hazır red sebepleri + özel mesaj).",
            "Mutfak için tarayıcıdan yazdır.",
            "Ödeme uygulama içinde alınmaz — müşteriyle telefonla netleşirsiniz.",
            "Teslimat sonrası müşteri Lezzet · Servis · Kurye puanı verebilir.",
        ],
    )

    _slide_content(
        prs,
        "Kötü yorumlar — Remedy (en büyük fark)",
        [
            "3 yıldız ve altı GastroSkor yorumları önce size özel gelir; herkese açık değil.",
            "24 saat içinde müşteriye indirim kuponu teklif edebilirsiniz (%5–100).",
            "Müşteri kabul ederse yorum gizli kalır; siz müşteriyi geri kazanırsınız.",
            "Süre dolarsa veya teklif vermezseniz yorum normal yayın akışına girer.",
            "Amaç: kamuya açık kötü puan yerine önce çözüm şansı.",
        ],
        footer="Panel → Düşük puanlı yorumlar / Remedy bölümü",
    )

    _slide_content(
        prs,
        "Özel şikayetler (Gizli geri bildirim)",
        [
            "Müşteriler mekan hakkında uygulama içi özel şikayet açabilir.",
            "Şikayetler herkese açık yorum değildir — sadece panelinizde görünür.",
            "Kategoriye göre filtreleyin; durum güncelleyin (inceleniyor / çözüldü).",
            "Müşteriye panelden mesaj yazın.",
            "Telafi kuponu tanımlayın (%10 / %25 / %50, süreli).",
        ],
        footer="Panel → Şikayetler",
    )

    _slide_content(
        prs,
        "Takipçiler ve kupon kampanyaları",
        [
            "Uygulamada mekanınızı takip eden kullanıcıları listede görün.",
            "Kampanya oluşturun: indirim %, geçerlilik günü, kupon adedi.",
            "Her takipçiye benzersiz GS- kupon kodu ve mobil bildirim gider.",
            "Kasada kodu girerek tek kullanımlık kuponu kapatabilirsiniz.",
            "Kampanya istatistikleri: verilen / kullanılan / süresi dolan.",
        ],
        footer="Panel → Takipçiler",
    )

    _slide_content(
        prs,
        "Rakip analizi (AI)",
        [
            "Rakip mekan ekleyin (denemede 1, ücretli planda daha fazla).",
            "AI, son ~90 gün yorumlarınızı rakiple kıyaslar.",
            "Rapor: sizin güçlü yönler, geliştirmeniz gereken alanlar, kanıt alıntıları.",
            "Geçmiş raporlar ve trend grafiği panelde saklanır.",
            "Standart planda belirli aralıklarla ücretsiz analiz hakkı; ek paket satın alınabilir.",
        ],
    )

    _slide_content(
        prs,
        "Google İşletme Profili",
        [
            "Google Business hesabınızı panele bağlayın (OAuth).",
            "Sadece 5 örnek yorum değil — tüm Google yorum havuzunu analiz edin.",
            "Trend ve tema özetleri panelde görünür.",
            "Bağlantıyı istediğiniz zaman kaldırabilirsiniz.",
        ],
    )

    _slide_content(
        prs,
        "Bildirimler",
        [
            "Panel zil ikonu: yeni sipariş, düşük puan, rakip değişimi, analiz hazır.",
            "E-posta tercihlerini açıp kapatabilirsiniz.",
            "Deneme bitişi ve sözleşme hatırlatmaları otomatik gelir.",
        ],
    )

    _slide_pricing(prs)

    _slide_content(
        prs,
        "İlk 7 gün — önerilen sıra",
        [
            "1. Panele girin ve mekan doğrulamasını tamamlayın.",
            "2. Menü + fotoğrafları yükleyin; online siparişi açın.",
            "3. Test siparişi alın; onay / red / yazdır akışını deneyin.",
            "4. Takipçi kampanyası veya remedy senaryosunu inceleyin.",
            "5. Bir rakip ekleyip AI raporu oluşturun.",
            "6. Sözleşmeyi imzalayıp posta ile gönderin (deneme bitmeden).",
        ],
        footer="Sorularınız için: paneldeki bildirimler ve başvuru e-postanızdaki iletişim kanalı.",
    )

    _slide_title(
        prs,
        "Hoş geldiniz!",
        "GastroSkor ailesine katıldığınız için teşekkürler.\n"
        "Panel: www.gastroskor.com.tr/panel",
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
